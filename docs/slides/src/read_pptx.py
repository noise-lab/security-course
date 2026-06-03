#!/usr/bin/env python3
"""
Script to read PowerPoint slides and extract text content
"""

try:
    from pptx import Presentation
    import sys

    pptx_file = "11-PrivacyLaw.pptx"

    prs = Presentation(pptx_file)

    print(f"Total slides: {len(prs.slides)}\n")
    print("="*80)

    for i, slide in enumerate(prs.slides, 1):
        print(f"\n### SLIDE {i}")
        print("-"*80)

        # Get slide title if exists
        if slide.shapes.title:
            print(f"Title: {slide.shapes.title.text}")
            print()

        # Get all text from shapes
        print("Content:")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                # Skip if it's the title (already printed)
                if slide.shapes.title and shape == slide.shapes.title:
                    continue
                print(f"  • {shape.text}")

        # Get speaker notes if they exist
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            if text_frame.text.strip():
                print("\nCurrent Speaker Notes:")
                print(f"  {text_frame.text}")

        print("="*80)

except ImportError:
    print("python-pptx not installed. Installing...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    print("\nPlease run the script again.")
    sys.exit(1)
except Exception as e:
    print(f"Error: {e}")
    sys.exit(1)
